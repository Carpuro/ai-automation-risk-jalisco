# -*- coding: utf-8 -*-
"""
Fill the university protocol template with the current thesis content.

Input : the CUCEA template (Protocolo de Investigacion - Presentacion .pptx)
Output: thesis/protocolo_avance_2026.pptx -- same design, updated title and
        filled body boxes (slides 2-6). The cronograma table (slide 7) is the
        student's official plan and is left untouched.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

# Template archived in-repo (the Downloads copy is disposable)
TEMPLATE = Path(__file__).resolve().parent / "plantilla_protocolo_cucea.pptx"
OUT = Path(__file__).resolve().parent / "protocolo_avance_2026.pptx"

TITLE = (
    "\n\nCarlos Pulido Rosas\n\n"
    "«Riesgo de automatización laboral por inteligencia artificial en "
    "Jalisco: exposición cognitiva y encarnada, incentivo económico y "
    "proyección 2025–2030»\n\n(caso Jalisco)"
)

BODIES = {
    2: [
        "La IA no es una tecnología: son DOS fronteras — la cognitiva (modelos "
        "de lenguaje) y la encarnada (robots, drones, máquinas autónomas) — con "
        "física, costos, canales de difusión y víctimas distintas.",
        "Jalisco las enfrenta a la vez: polo de servicios modernos "
        "(Guadalajara tech) y corredor industrial-agroalimentario que el "
        "nearshoring expande.",
        "El 55% del empleo es informal: la misma exposición, sin red alguna "
        "de protección (sin IMSS, sin indemnización).",
        "El debate público pregunta «¿cuántos empleos destruirá la IA?» — una "
        "pregunta mal planteada que no distingue frontera, incentivo "
        "económico ni protección social.",
    ],
    3: [
        "Pregunta de investigación: ¿cómo se distribuye el riesgo de "
        "automatización entre la frontera cognitiva y la encarnada, y qué "
        "determina cuál domina en cada ocupación, sector y territorio?",
        "H1 — Gradiente: la exposición a LLMs sigue la jerarquía ocupacional "
        "(invierte la «zona segura» de Frey-Osborne).",
        "H2 — Estructura bipolar: las dos exposiciones son polos opuestos de "
        "una dimensión dominante, no riesgos independientes.",
        "H3 — Posición de Jalisco: por debajo del centro cognitivo y por "
        "encima del encarnado.",
        "H4 — Moderación económica: la exposición solo se vuelve presión real "
        "donde automatizar es RENTABLE (salario / costo de capital).",
    ],
    4: [
        "Marco de tareas (Autor-Levy-Murnane 2003; Acemoglu-Autor 2011): la "
        "tecnología sustituye tareas, no empleos → la exposición se mide del "
        "contenido ocupacional (O*NET).",
        "Carrera desplazamiento vs. reinstalación + condición de rentabilidad "
        "(Acemoglu-Restrepo 2018-19): factibilidad ≠ adopción → el IRA y la "
        "hipótesis H4; «tecnologías mediocres» = crecimiento perdido.",
        "Paradoja de Moravec: dos fronteras con costos opuestos → predicción "
        "de exposiciones OPUESTAS (H2).",
        "Linaje de medición: Frey-Osborne → Webb (patentes) → Felten (AIOE) → "
        "Eloundou (GPT) → Índice Económico de Anthropic (uso observado). "
        "Contribución propia: DBOE (dinamiza a Felten con benchmarks "
        "medidos) + DEOE (el espejo encarnado que faltaba).",
        "Curva J de la productividad (Brynjolfsson): el rezago entre choque "
        "de capacidad y choque laboral → la VENTANA DE POLÍTICA.",
        "Economías en desarrollo: informalidad como multiplicador de "
        "severidad; trasplante O*NET contrastado (r = 0.97 MX-EUA); "
        "nearshoring como importador de la frontera robótica.",
    ],
    5: [
        "14 fuentes integradas en SQL Server (~45 tablas); 100% reproducible: "
        "cada cifra proviene de un script versionado (GitHub).",
        "Índices propios validados: DBOE reproduce el AIOE publicado "
        "(r = 0.94); DEOE valida contra exposición robótica por patentes de "
        "Webb (r = +0.76) con validez discriminante.",
        "Cruce oficial INEGI SINCO↔SOC: 88% del empleo ponderado de Jalisco "
        "con exposición asignada; análisis de cotas para el resto.",
        "Nivel 1 — modelo hurdle del USO OBSERVADO de IA por ocupación "
        "(objetivo externo, no circular).",
        "Contraste de H4 — panel censal 2003–2023 (19 sectores × 4 "
        "transiciones), IRA rezagado, inferencia exacta por permutación.",
        "Nivel 2 — escenarios 2025–2030: presión = exposición × curvas "
        "tecnológicas c(t)/r(t) × incentivo (IRA).",
    ],
    6: [
        "H2 ✓ Eje bipolar (57% de varianza; φ = −0.67): los LLMs y los robots "
        "amenazan ocupaciones OPUESTAS — y los 32 estados del país se alinean "
        "sobre esa misma recta.",
        "H1 ✓ Gradiente espejo por educación: DBOE de −1.00 (primaria) a "
        "+0.50 (superior); DEOE exactamente al revés. El riesgo tiene género: "
        "mujeres → LLMs; hombres → robots.",
        "H3 ✓ refinada: Jalisco es BIFRONTAL — 8.º estado más expuesto a LLMs "
        "y 29.º a robots; metrópoli cognitiva vs. corredor encarnado.",
        "H4 ✓ Veinte años de censos: el mercado sustituye donde paga "
        "(profundización de capital p = .013; participación laboral p = .004); "
        "la forma del daño es crecimiento perdido: +0.6 vs. +3.7% anual.",
        "Cuantificación: 1.10 M de trabajadores bajo presión alta hoy → 1.88 M "
        "en 2030 (escenario base); 908 mil informales sin red; celda crítica "
        "de 348 mil (informal + robot + ≤9 años de escolaridad).",
        "La ventana de política: el choque de capacidad ya ocurrió y las "
        "expectativas saltaron +12.5 pp (2020→2023), pero el empleo —formal, "
        "total o vía informalidad— aún no responde. La ventana sigue abierta.",
    ],
}


def fill_box(shape, bullets, size=11):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, b in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = "• " + b
        para.space_after = Pt(6)
        for run in para.runs:
            run.font.size = Pt(size)


def main():
    p = Presentation(TEMPLATE)
    slides = list(p.slides)

    # Slide 1: new title; date box -> avance 2026
    slides[0].shapes[4].text_frame.text = TITLE
    for run_para in slides[0].shapes[4].text_frame.paragraphs:
        for r in run_para.runs:
            r.font.size = Pt(16)
    slides[0].shapes[1].text_frame.text = "Avance de tesis — junio de 2026"

    # Slides 2-6: fill the empty body boxes (shape index 0)
    for idx, bullets in BODIES.items():
        fill_box(slides[idx - 1].shapes[0], bullets)

    p.save(OUT)
    print(f"-> {OUT}")


if __name__ == "__main__":
    main()
